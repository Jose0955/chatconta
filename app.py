import streamlit as st
from groq import Groq
from datetime import datetime
import pandas as pd
import io
import os
import re
import json
import time
import random
import streamlit.components.v1 as components
from openpyxl.styles import Font


def hojas_excel_a_markdown(hojas: dict) -> str:
    """Convierte un diccionario de hojas (nombre -> DataFrame) de un Excel
    en texto con tablas Markdown. Función compartida: la usan tanto el
    material que suben los docentes como los ejercicios que suben los
    estudiantes."""
    partes = []
    for nombre_hoja, df in hojas.items():
        partes.append(f"--- Hoja: {nombre_hoja} ---")
        partes.append(df.to_markdown(index=False))
    return "\n".join(partes)


# ---------------------------------------------------------
# MATERIAL DE CLASE SUBIDO POR LOS DOCENTES
# Los docentes NO suben archivos desde la app (eso requeriría una base de
# datos); en vez de eso, suben sus archivos directamente a la carpeta
# "material_docentes/" del repositorio de GitHub, y Contín los lee solo
# al arrancar. Soporta: .txt, .md, .xlsx/.xls, .pdf, .pptx
# Esto se carga AL INICIO del archivo (antes que todo lo demás) porque
# tanto la barra lateral como el prompt del sistema lo necesitan.
# ---------------------------------------------------------
CARPETA_MATERIAL_DOCENTES = "material_docentes"
LARGO_MAXIMO_MATERIAL = 6000  # límite de caracteres para no disparar el consumo de tokens


@st.cache_data(show_spinner=False)
def cargar_material_docentes():
    """Lee todos los archivos de la carpeta material_docentes/ y arma un
    texto resumido para dárselo a Contín como referencia extra. Si un
    archivo falla al leerse, simplemente se lo salta (no rompe la app)."""
    if not os.path.isdir(CARPETA_MATERIAL_DOCENTES):
        return "", []

    partes = []
    archivos_cargados = []

    for nombre_archivo in sorted(os.listdir(CARPETA_MATERIAL_DOCENTES)):
        if nombre_archivo.startswith(".") or nombre_archivo.upper().startswith("README"):
            continue
        ruta = os.path.join(CARPETA_MATERIAL_DOCENTES, nombre_archivo)
        if not os.path.isfile(ruta):
            continue
        extension = nombre_archivo.lower().rsplit(".", 1)[-1] if "." in nombre_archivo else ""

        try:
            texto_archivo = None

            if extension in ("txt", "md"):
                with open(ruta, "r", encoding="utf-8", errors="ignore") as f:
                    texto_archivo = f.read()

            elif extension in ("xlsx", "xls"):
                hojas = pd.read_excel(ruta, sheet_name=None)
                texto_archivo = hojas_excel_a_markdown(hojas)

            elif extension == "pdf":
                from pypdf import PdfReader
                lector = PdfReader(ruta)
                paginas_texto = [(p.extract_text() or "") for p in lector.pages[:20]]
                texto_archivo = "\n".join(paginas_texto)

            elif extension == "pptx":
                from pptx import Presentation
                presentacion = Presentation(ruta)
                lineas = []
                for diapositiva in presentacion.slides:
                    for forma in diapositiva.shapes:
                        if forma.has_text_frame:
                            texto_forma = forma.text_frame.text.strip()
                            if texto_forma:
                                lineas.append(texto_forma)
                texto_archivo = "\n".join(lineas)

            if texto_archivo and texto_archivo.strip():
                partes.append(f"--- Material del docente: {nombre_archivo} ---\n{texto_archivo.strip()}")
                archivos_cargados.append(nombre_archivo)

        except Exception:
            # Si un archivo específico falla (formato raro, corrupto, etc.),
            # lo saltamos sin tumbar el resto de la app.
            continue

    texto_completo = "\n\n".join(partes)
    if len(texto_completo) > LARGO_MAXIMO_MATERIAL:
        texto_completo = texto_completo[:LARGO_MAXIMO_MATERIAL] + "\n[...material recortado por espacio...]"

    return texto_completo, archivos_cargados


MATERIAL_DOCENTES_TEXTO, MATERIAL_DOCENTES_ARCHIVOS = cargar_material_docentes()

# =========================================================
# CONFIGURACIÓN GENERAL DE LA PÁGINA
# =========================================================
st.set_page_config(
    page_title="Contín - Tu Tutor de Contabilidad",
    page_icon="🤝",
    layout="centered"
)

# ---------------------------------------------------------
# ESTILOS PERSONALIZADOS — modo oscuro fijo, estilo "Gmail oscuro"
# (fondo casi negro + acentos azules), con estrellitas animadas.
# Al ser un tema FIJO (no depende del modo claro/oscuro del celular
# o la laptop), el contraste de texto siempre queda correcto.
# ---------------------------------------------------------
AZUL_ACENTO = "#8AB4F8"     # azul estilo Google/Gmail modo oscuro
AZUL_SUAVE = "#5B9BF0"
FONDO_APP = "#050709"        # negro más duro y sólido
FONDO_TARJETA = "#12161C"
FONDO_BURBUJA_USUARIO = "#1B2026"
FONDO_BURBUJA_ASISTENTE = "#0F1B2A"
TEXTO_CLARO = "#FFFFFF"

# Generamos posiciones aleatorias (pero fijas por sesión) para las estrellitas
random.seed(7)
_estrellas_html = ""
for i in range(35):
    top = random.uniform(0, 100)
    left = random.uniform(0, 100)
    tamano = random.uniform(2, 4)
    duracion = random.uniform(4, 9)
    retraso = random.uniform(0, 6)
    _estrellas_html += (
        f'<div class="estrella" style="'
        f'top:{top}%; left:{left}%; width:{tamano}px; height:{tamano}px; '
        f'animation-duration:{duracion}s; animation-delay:{retraso}s;"></div>'
    )

st.markdown(
    f"""
    <style>
    /* ---------- Fondo general y estrellitas animadas ---------- */
    .stApp {{
        background: {FONDO_APP} !important;
    }}
    html, body {{
        background-color: {FONDO_APP} !important;
    }}
    /* La zona de abajo donde vive la caja de texto (queda fuera de .stApp
       en algunas versiones de Streamlit y se veía blanca) */
    [data-testid="stBottom"],
    [data-testid="stBottomBlockContainer"],
    .stChatFloatingInputContainer,
    [data-testid="stAppViewContainer"] {{
        background-color: {FONDO_APP} !important;
    }}
    /* Barra decorativa roja de arriba -> la pasamos a azul para que combine */
    [data-testid="stDecoration"] {{
        background-image: linear-gradient(90deg, {AZUL_SUAVE}, {AZUL_ACENTO}) !important;
    }}

    .campo-estrellas {{
        position: fixed;
        top: 0; left: 0;
        width: 100%; height: 100%;
        overflow: hidden;
        pointer-events: none;
        z-index: 0;
    }}
    .estrella {{
        position: absolute;
        background: {AZUL_ACENTO};
        border-radius: 50%;
        opacity: 0.25;
        box-shadow: 0 0 6px 1px {AZUL_ACENTO};
        animation-name: flotar, titilar;
        animation-iteration-count: infinite;
        animation-timing-function: ease-in-out;
    }}
    @keyframes flotar {{
        0%   {{ transform: translateY(0px); }}
        50%  {{ transform: translateY(-18px); }}
        100% {{ transform: translateY(0px); }}
    }}
    @keyframes titilar {{
        0%, 100% {{ opacity: 0.15; }}
        50%      {{ opacity: 0.7; }}
    }}

    /* Todo el contenido va por encima del campo de estrellas */
    .main .block-container {{
        padding-top: 2rem;
        position: relative;
        z-index: 1;
    }}

    /* ---------- Texto general (arregla el problema de contraste) ---------- */
    .stApp, .stApp p, .stApp li, .stApp span, .stApp label,
    [data-testid="stMarkdownContainer"] {{
        color: {TEXTO_CLARO} !important;
    }}
    h1 {{
        color: {AZUL_ACENTO} !important;
    }}
    h2, h3, h4 {{
        color: {AZUL_SUAVE} !important;
    }}
    a {{ color: {AZUL_ACENTO} !important; }}

    /* ---------- Barra lateral ---------- */
    section[data-testid="stSidebar"] {{
        background-color: {FONDO_TARJETA} !important;
        border-right: 1px solid #23293380;
    }}
    section[data-testid="stSidebar"] * {{
        color: {TEXTO_CLARO} !important;
    }}

    /* ---------- Burbujas de chat ---------- */
    [data-testid="stChatMessage"] {{
        border-radius: 16px;
        border: 1px solid #23293380;
    }}
    [data-testid="stChatMessage"]:nth-of-type(odd) {{
        background-color: {FONDO_BURBUJA_USUARIO} !important;
    }}
    [data-testid="stChatMessage"]:nth-of-type(even) {{
        background-color: {FONDO_BURBUJA_ASISTENTE} !important;
        border-left: 3px solid {AZUL_ACENTO};
    }}

    /* ---------- Caja de texto del chat ---------- */
    [data-testid="stChatInput"] textarea {{
        background-color: {FONDO_TARJETA} !important;
        color: {TEXTO_CLARO} !important;
    }}
    [data-testid="stChatInput"] {{
        background-color: {FONDO_TARJETA} !important;
        border: 1px solid {AZUL_ACENTO}55 !important;
    }}

    /* ---------- Botones ---------- */
    .stButton button, .stDownloadButton button {{
        border-radius: 12px;
        background-color: {AZUL_ACENTO} !important;
        color: #0B0E14 !important;
        border: none !important;
        font-weight: 600;
    }}
    .stButton button:hover, .stDownloadButton button:hover {{
        background-color: {AZUL_SUAVE} !important;
    }}

    /* ---------- Expanders (paneles de voz y Excel) ---------- */
    [data-testid="stExpander"] {{
        background-color: {FONDO_TARJETA} !important;
        border-radius: 12px;
        border: 1px solid #23293380;
    }}

    /* ---------- Radios (nivel) ---------- */
    [data-testid="stSidebar"] [role="radiogroup"] label {{
        color: {TEXTO_CLARO} !important;
    }}
    </style>

    <div class="campo-estrellas">{_estrellas_html}</div>
    """,
    unsafe_allow_html=True,
)

st.title("🤝 Contín, tu asistente contable de confianza")

# ---------------------------------------------------------
# MASCOTA: Contín, el alien-pulpo contable 🐙
# Cambia de cara según el momento: pensando, hablando o feliz.
# ---------------------------------------------------------
def mascota_svg(estado: str = "normal", modo_hero: bool = False) -> str:
    if estado == "pensando":
        ojos = """
            <circle cx="78" cy="95" r="13" fill="white"/>
            <circle cx="122" cy="95" r="13" fill="white"/>
            <circle cx="83" cy="90" r="6" fill="#0B3D91"/>
            <circle cx="127" cy="90" r="6" fill="#0B3D91"/>
        """
        boca = '<ellipse cx="100" cy="128" rx="7" ry="6" fill="#0B3D91"/>'
        extra = """
            <g class="burbuja-pensar">
                <circle cx="150" cy="55" r="5" fill="white" opacity="0.85"/>
                <circle cx="163" cy="42" r="7" fill="white" opacity="0.85"/>
                <ellipse cx="182" cy="24" rx="16" ry="12" fill="white" opacity="0.9"/>
                <text x="182" y="29" font-size="14" text-anchor="middle" fill="#5B9BF0">?</text>
            </g>
        """
    elif estado == "hablando":
        ojos = """
            <circle cx="78" cy="95" r="14" fill="white"/>
            <circle cx="122" cy="95" r="14" fill="white"/>
            <circle cx="80" cy="95" r="7" fill="#0B3D91"/>
            <circle cx="124" cy="95" r="7" fill="#0B3D91"/>
        """
        boca = '<ellipse cx="100" cy="130" rx="14" ry="11" fill="#0B3D91"/>'
        extra = """
            <g class="chispas">
                <path d="M158 60 L162 70 L172 72 L162 76 L158 86 L154 76 L144 72 L154 70 Z" fill="white" opacity="0.9"/>
                <circle cx="35" cy="65" r="4" fill="white" opacity="0.7"/>
            </g>
        """
    elif estado == "feliz":
        ojos = """
            <path d="M68 95 Q78 82 88 95" stroke="white" stroke-width="6" fill="none" stroke-linecap="round"/>
            <path d="M112 95 Q122 82 132 95" stroke="white" stroke-width="6" fill="none" stroke-linecap="round"/>
        """
        boca = '<path d="M78 122 Q100 148 122 122" stroke="#0B3D91" stroke-width="7" fill="none" stroke-linecap="round"/>'
        extra = """
            <g class="corazones">
                <text x="35" y="55" font-size="22">💙</text>
                <text x="160" y="45" font-size="18">✨</text>
                <text x="150" y="90" font-size="16">💙</text>
            </g>
        """
    elif estado == "bailando":
        ojos = """
            <path d="M68 95 Q78 82 88 95" stroke="white" stroke-width="6" fill="none" stroke-linecap="round"/>
            <path d="M112 95 Q122 82 132 95" stroke="white" stroke-width="6" fill="none" stroke-linecap="round"/>
        """
        boca = '<path d="M78 122 Q100 148 122 122" stroke="#0B3D91" stroke-width="7" fill="none" stroke-linecap="round"/>'
        extra = """
            <g class="notas-musicales">
                <text x="30" y="50" font-size="22">🎵</text>
                <text x="158" y="40" font-size="20">🎶</text>
            </g>
        """
    elif estado == "cantando":
        ojos = """
            <path d="M68 92 Q78 80 88 92" stroke="white" stroke-width="6" fill="none" stroke-linecap="round"/>
            <path d="M112 92 Q122 80 132 92" stroke="white" stroke-width="6" fill="none" stroke-linecap="round"/>
        """
        boca = '<ellipse cx="100" cy="132" rx="12" ry="14" fill="#0B3D91"/>'
        extra = """
            <g class="notas-musicales">
                <text x="150" y="45" font-size="24">🎵</text>
                <text x="28" y="60" font-size="18">🎶</text>
            </g>
            <g class="microfono">
                <rect x="150" y="98" width="12" height="26" rx="6" fill="#EEE"/>
                <line x1="156" y1="124" x2="156" y2="145" stroke="#5B9BF0" stroke-width="4"/>
                <circle cx="156" cy="96" r="10" fill="#DDD"/>
            </g>
        """
    else:  # normal / idle
        ojos = """
            <circle cx="78" cy="95" r="13" fill="white" class="parpadeo"/>
            <circle cx="122" cy="95" r="13" fill="white" class="parpadeo"/>
            <circle cx="78" cy="95" r="6" fill="#0B3D91"/>
            <circle cx="122" cy="95" r="6" fill="#0B3D91"/>
        """
        boca = '<path d="M85 125 Q100 135 115 125" stroke="#0B3D91" stroke-width="5" fill="none" stroke-linecap="round"/>'
        extra = ""

    # ---------------------------------------------------------
    # Accesorios: en modo "profesional" (explicando conta) usa
    # lentes + calculadora. En modo "amigo" (consejos, baile,
    # canto, celebración) se los quita y quedan tirados al lado.
    # ---------------------------------------------------------
    ESTADOS_PROFESIONALES = ("normal", "pensando", "hablando")
    if estado in ESTADOS_PROFESIONALES:
        accesorios = """
            <g class="lentes">
                <circle cx="78" cy="95" r="18" fill="#8AB4F8" fill-opacity="0.25" stroke="#0B3D91" stroke-width="3"/>
                <circle cx="122" cy="95" r="18" fill="#8AB4F8" fill-opacity="0.25" stroke="#0B3D91" stroke-width="3"/>
                <line x1="96" y1="95" x2="104" y2="95" stroke="#0B3D91" stroke-width="3"/>
            </g>
            <g class="calculadora">
                <rect x="150" y="118" width="26" height="34" rx="4" fill="#E8EAED" stroke="#3E7BD9" stroke-width="2"/>
                <rect x="154" y="122" width="18" height="7" rx="1" fill="#3E7BD9"/>
                <circle cx="157" cy="135" r="2" fill="#3E7BD9"/>
                <circle cx="163" cy="135" r="2" fill="#3E7BD9"/>
                <circle cx="169" cy="135" r="2" fill="#3E7BD9"/>
                <circle cx="157" cy="143" r="2" fill="#3E7BD9"/>
                <circle cx="163" cy="143" r="2" fill="#3E7BD9"/>
                <circle cx="169" cy="143" r="2" fill="#3E7BD9"/>
            </g>
        """
    else:
        # Modo amigo: lentes y calculadora tirados a un lado
        accesorios = """
            <g class="modo-amigo-doodle" opacity="0.8">
                <circle cx="182" cy="168" r="7" fill="none" stroke="#5B9BF0" stroke-width="2"/>
                <circle cx="196" cy="172" r="7" fill="none" stroke="#5B9BF0" stroke-width="2"/>
                <line x1="189" y1="169" x2="189" y2="171" stroke="#5B9BF0" stroke-width="2"/>
                <rect x="8" y="172" width="16" height="20" rx="3" fill="#E8EAED" stroke="#3E7BD9" stroke-width="2" transform="rotate(-18 16 182)"/>
            </g>
        """

    clase_extra = " mascota-bailando" if estado == "bailando" else ""
    clase_extra += " mascota-hero" if modo_hero else ""

    svg_html = f"""
    <div class="mascota-flotante{clase_extra}">
    <svg viewBox="0 0 200 220" width="150" height="165" xmlns="http://www.w3.org/2000/svg">
        <defs>
            <radialGradient id="cuerpoGrad" cx="40%" cy="35%" r="75%">
                <stop offset="0%" stop-color="#EAF3FF"/>
                <stop offset="45%" stop-color="#8AB4F8"/>
                <stop offset="100%" stop-color="#3E7BD9"/>
            </radialGradient>
        </defs>
        <g class="tentaculo t1"><path d="M55 165 Q40 190 50 215" stroke="#5B9BF0" stroke-width="14" fill="none" stroke-linecap="round"/></g>
        <g class="tentaculo t2"><path d="M80 175 Q75 200 82 218" stroke="#5B9BF0" stroke-width="14" fill="none" stroke-linecap="round"/></g>
        <g class="tentaculo t3"><path d="M120 175 Q125 200 118 218" stroke="#5B9BF0" stroke-width="14" fill="none" stroke-linecap="round"/></g>
        <g class="tentaculo t4"><path d="M145 165 Q160 190 150 215" stroke="#5B9BF0" stroke-width="14" fill="none" stroke-linecap="round"/></g>
        <line x1="70" y1="45" x2="55" y2="15" stroke="#5B9BF0" stroke-width="5" stroke-linecap="round"/>
        <circle cx="55" cy="12" r="7" fill="#EAF3FF" class="antena"/>
        <line x1="130" y1="45" x2="145" y2="15" stroke="#5B9BF0" stroke-width="5" stroke-linecap="round"/>
        <circle cx="145" cy="12" r="7" fill="#EAF3FF" class="antena"/>
        <ellipse cx="100" cy="110" rx="80" ry="75" fill="url(#cuerpoGrad)"/>
        {ojos}
        {boca}
        {extra}
        {accesorios}
    </svg>
    </div>
    """
    # Importante: quitamos cualquier línea en blanco, porque Streamlit
    # (al interpretar esto como Markdown) corta el bloque de HTML apenas
    # encuentra una línea vacía, y el resto se muestra como texto crudo.
    return "\n".join(linea for linea in svg_html.split("\n") if linea.strip() != "")


def lanzar_confeti():
    """Dispara una animación de confeti de colores. Usa canvas-confetti dentro
    de un componente HTML (necesario porque Streamlit no ejecuta <script>
    sueltos en st.markdown), con un tamaño real y visible de entrada, e
    intenta además expandirse a toda la pantalla."""
    components.html(
        """
        <script src="https://cdn.jsdelivr.net/npm/canvas-confetti@1.9.3/dist/confetti.browser.min.js"></script>
        <script>
        (function () {
            try {
                var marco = window.frameElement;
                if (marco) {
                    marco.style.position = 'fixed';
                    marco.style.top = '0';
                    marco.style.left = '0';
                    marco.style.width = '100vw';
                    marco.style.height = '100vh';
                    marco.style.zIndex = '999999';
                    marco.style.pointerEvents = 'none';
                    marco.style.border = 'none';
                }
            } catch (e) {}

            function empezar() {
                if (typeof confetti !== 'function') { setTimeout(empezar, 100); return; }
                var colores = ['#8AB4F8', '#5B9BF0', '#FFFFFF', '#EAF3FF', '#FFD166', '#FF6B6B'];
                confetti({ particleCount: 200, spread: 120, origin: { y: 0.3 }, colors: colores });
                var fin = Date.now() + 2800;
                (function ciclo() {
                    confetti({ particleCount: 8, angle: 60, spread: 80, origin: { x: 0, y: 0.6 }, colors: colores });
                    confetti({ particleCount: 8, angle: 120, spread: 80, origin: { x: 1, y: 0.6 }, colors: colores });
                    if (Date.now() < fin) { requestAnimationFrame(ciclo); }
                })();
            }
            empezar();
        })();
        </script>
        """,
        height=500,
    )
    # Respaldo garantizado: el efecto nativo de celebración de Streamlit,
    # que SIEMPRE funciona (no depende de scripts externos ni de CDNs).
    st.balloons()


def escribir_con_efecto_maquina(texto: str, placeholder=None):
    """Muestra el texto poco a poco, como si Contín lo estuviera escribiendo
    en vivo. Si el texto es muy largo, acelera para no hacer esperar de más."""
    if placeholder is None:
        placeholder = st.empty()

    palabras = texto.split(" ")
    # Textos largos (muchas palabras) se revelan más rápido para no ser pesados
    velocidad = 0.028 if len(palabras) < 40 else (0.014 if len(palabras) < 120 else 0.006)

    acumulado = ""
    for palabra in palabras:
        acumulado += palabra + " "
        placeholder.markdown(acumulado + "▌")
        time.sleep(velocidad)
    placeholder.markdown(acumulado.strip())


def calcular_van(tasa: float, flujos: list) -> float:
    """Calcula el Valor Actual Neto. 'flujos[0]' debe ser la inversión
    inicial (negativa), y el resto los flujos de caja de cada periodo."""
    return sum(flujo / (1 + tasa) ** i for i, flujo in enumerate(flujos))


def calcular_tir(flujos: list):
    """Calcula la Tasa Interna de Retorno por bisección (búsqueda binaria).
    Devuelve None si no encuentra una tasa razonable entre -99% y 1000%."""
    baja, alta = -0.99, 10.0
    van_baja = calcular_van(baja, flujos)
    van_alta = calcular_van(alta, flujos)
    if van_baja * van_alta > 0:
        return None
    for _ in range(200):
        medio = (baja + alta) / 2
        van_medio = calcular_van(medio, flujos)
        if abs(van_medio) < 0.01:
            return medio
        if van_baja * van_medio < 0:
            alta = medio
            van_alta = van_medio
        else:
            baja = medio
            van_baja = van_medio
    return (baja + alta) / 2


def limpiar_para_voz(texto: str) -> str:
    """Prepara el texto de Contín para leerlo en voz alta: quita símbolos de
    Markdown (asteriscos, gatos, barras, etc.) y reemplaza las tablas por una
    frase corta, para que no suene raro al escucharlo."""
    lineas = texto.split("\n")
    resultado = []
    dentro_tabla = False
    for linea in lineas:
        if linea.strip().startswith("|"):
            if not dentro_tabla:
                resultado.append("Te dejé una tabla en pantalla con el detalle completo.")
                dentro_tabla = True
            continue
        dentro_tabla = False
        resultado.append(linea)

    texto_limpio = "\n".join(resultado)
    texto_limpio = re.sub(r'[*_#`>]', '', texto_limpio)
    texto_limpio = re.sub(r'\n{2,}', '. ', texto_limpio)
    texto_limpio = re.sub(r'\s+', ' ', texto_limpio).strip()
    return texto_limpio


def hablar_texto(texto: str):
    """Hace que el navegador lea el texto en voz alta (texto-a-voz nativo,
    no necesita ningún servicio externo, funciona sin internet)."""
    texto_para_hablar = limpiar_para_voz(texto)
    if not texto_para_hablar:
        return
    texto_js = json.dumps(texto_para_hablar)
    components.html(
        f"""
        <script>
        try {{
            window.speechSynthesis.cancel();
            var utterance = new SpeechSynthesisUtterance({texto_js});
            utterance.lang = 'es-ES';
            utterance.rate = 1.02;
            utterance.pitch = 1.05;
            window.speechSynthesis.speak(utterance);
        }} catch (e) {{}}
        </script>
        """,
        height=0,
    )


# ---------------------------------------------------------
# QUIZZES AUTOMÁTICOS 🎯
# ---------------------------------------------------------
PROMPT_QUIZ_SISTEMA = """Eres un generador de quizzes de opción múltiple sobre
contabilidad para estudiantes de Bachillerato Técnico en Ecuador.

Devuelve ÚNICAMENTE un JSON válido (sin texto adicional, sin explicaciones,
sin marcadores de código como ```), con esta forma EXACTA:

[
  {"pregunta": "...", "opciones": ["...", "...", "...", "..."], "respuesta_correcta": 0, "explicacion": "..."}
]

Reglas:
- TÚ decides cuántas preguntas hacer (entre 3 y 6) según qué tan amplio sea
  el tema: un tema puntual merece 3 preguntas, un tema amplio puede llegar a 6.
- Cada pregunta debe tener EXACTAMENTE 4 opciones.
- "respuesta_correcta" es el índice (0, 1, 2 o 3) de la opción correcta.
- Las preguntas deben basarse específicamente en el tema/contexto que te den,
  no en contabilidad en general.
- "explicacion" es una frase corta (máx. 2 líneas) de por qué esa es la
  respuesta correcta.
- No agregues nada fuera del JSON."""


def generar_quiz(tema_contexto: str):
    """Le pide a la IA un quiz en formato JSON sobre el tema dado.
    Devuelve una lista de preguntas, o None si algo falla."""
    try:
        respuesta = client.chat.completions.create(
            model=MODEL_NAME,
            messages=[
                {"role": "system", "content": PROMPT_QUIZ_SISTEMA},
                {"role": "user", "content": f"Tema/contexto de la clase:\n{tema_contexto}\n\nGenera el quiz."},
            ],
        )
        texto = respuesta.choices[0].message.content.strip()
        # Por si el modelo igual mete ``` alrededor del JSON, lo limpiamos
        texto = re.sub(r"^```(json)?|```$", "", texto.strip(), flags=re.MULTILINE).strip()

        preguntas = json.loads(texto)
        preguntas_validas = []
        for p in preguntas:
            if (
                isinstance(p, dict)
                and "pregunta" in p and "opciones" in p and "respuesta_correcta" in p
                and len(p["opciones"]) >= 2
                and 0 <= p["respuesta_correcta"] < len(p["opciones"])
            ):
                preguntas_validas.append(p)

        return preguntas_validas if preguntas_validas else None
    except Exception as e:
        st.error(f"No se pudo generar el quiz. Detalle técnico: {e}")
        return None


def boton_generar_quiz(texto_contexto: str, key_sufijo: str):
    """Botón reutilizable que aparece después de una respuesta o en la
    barra lateral, para pedirle a Contín un quiz sobre ese tema."""
    if st.button("🎯 Hazme un quiz de esto", key=f"quiz_btn_{key_sufijo}"):
        with st.spinner("Contín está armando tu quiz..."):
            preguntas = generar_quiz(texto_contexto)
        if preguntas:
            st.session_state.quiz_id = st.session_state.get("quiz_id", 0) + 1
            st.session_state.quiz_actual = {"preguntas": preguntas, "tema": texto_contexto[:80]}
            st.rerun()


def boton_explicar_mas_facil(key_sufijo: str):
    """Botón que le pide a Contín que reexplique su última respuesta de
    forma más sencilla, sin que el estudiante tenga que reescribir nada."""
    if st.button("🔁 Explícamelo más fácil", key=f"facil_btn_{key_sufijo}"):
        responder_pregunta(
            "Por favor, explícame lo mismo de tu respuesta anterior pero de una "
            "forma más fácil de entender: usa palabras más sencillas, ve más "
            "despacio paso a paso, y si puedes, dame un ejemplo distinto al anterior."
        )


def mostrar_quiz():
    """Dibuja el quiz activo (si hay uno) como un formulario interactivo."""
    quiz = st.session_state.get("quiz_actual")
    if not quiz:
        return

    preguntas = quiz["preguntas"]
    qid = st.session_state.get("quiz_id", 0)

    st.markdown("### 🎯 Quiz rápido")
    with st.form(key=f"form_quiz_{qid}"):
        seleccionadas = []
        for i, p in enumerate(preguntas):
            opciones_texto = [f"{chr(65 + j)}. {op}" for j, op in enumerate(p["opciones"])]
            seleccion = st.radio(
                f"**{i + 1}. {p['pregunta']}**",
                opciones_texto,
                key=f"quiz_{qid}_p{i}",
                index=None,
            )
            seleccionadas.append(seleccion)
        enviado = st.form_submit_button("✅ Revisar respuestas")

    if enviado:
        correctas = 0
        for i, p in enumerate(preguntas):
            seleccion = seleccionadas[i]
            idx_elegido = (ord(seleccion[0]) - 65) if seleccion else -1
            es_correcta = idx_elegido == p["respuesta_correcta"]
            if es_correcta:
                correctas += 1
            texto_correcta = p["opciones"][p["respuesta_correcta"]]
            if es_correcta:
                st.success(f"**{i + 1}.** ¡Correcto! {p.get('explicacion', '')}")
            else:
                st.error(f"**{i + 1}.** La respuesta correcta era: {texto_correcta}. {p.get('explicacion', '')}")

        st.markdown(f"## Resultado: {correctas}/{len(preguntas)}")

        # ---- Racha y logros (viven solo en esta sesión del navegador) ----
        if "racha_quiz" not in st.session_state:
            st.session_state.racha_quiz = 0
        if "total_quizzes_perfectos" not in st.session_state:
            st.session_state.total_quizzes_perfectos = 0
        if "logros_desbloqueados" not in st.session_state:
            st.session_state.logros_desbloqueados = set()

        if correctas == len(preguntas):
            st.session_state.mascota_estado = "feliz"
            st.session_state.racha_quiz += 1
            st.session_state.total_quizzes_perfectos += 1
            lanzar_confeti()

            nuevos_logros = []
            hitos = {
                1: "🥇 ¡Primer quiz perfecto!",
                3: "🔥 3 quizzes perfectos seguidos",
                5: "🌟 5 quizzes perfectos seguidos",
                10: "🏆 10 quizzes perfectos seguidos",
            }
            for hito, texto_logro in hitos.items():
                if st.session_state.racha_quiz == hito and texto_logro not in st.session_state.logros_desbloqueados:
                    st.session_state.logros_desbloqueados.add(texto_logro)
                    nuevos_logros.append(texto_logro)
            for logro in nuevos_logros:
                st.toast(logro, icon="🏆")
        else:
            st.session_state.racha_quiz = 0  # se rompe la racha si no fue perfecto

    if st.button("✖️ Cerrar quiz"):
        st.session_state.quiz_actual = None
        st.rerun()


if "mascota_estado" not in st.session_state:
    st.session_state.mascota_estado = "normal"

mascota_placeholder = st.empty()
_estado_inicial = "bailando" if st.session_state.get("bailando", False) else st.session_state.mascota_estado
_es_primera_visita = len(st.session_state.get("messages", [])) == 0
with mascota_placeholder.container():
    st.markdown(mascota_svg(_estado_inicial, modo_hero=_es_primera_visita), unsafe_allow_html=True)

st.markdown(
    f"""
    <style>
    .tentaculo {{ transform-origin: top center; animation: ondear 3s ease-in-out infinite; }}
    .t1 {{ animation-delay: 0s; }} .t2 {{ animation-delay: 0.4s; }}
    .t3 {{ animation-delay: 0.2s; }} .t4 {{ animation-delay: 0.6s; }}
    @keyframes ondear {{
        0%, 100% {{ transform: rotate(0deg); }}
        50% {{ transform: rotate(6deg); }}
    }}
    .antena {{ animation: brillo 2s ease-in-out infinite; }}
    @keyframes brillo {{
        0%, 100% {{ opacity: 0.6; filter: drop-shadow(0 0 2px {AZUL_ACENTO}); }}
        50% {{ opacity: 1; filter: drop-shadow(0 0 8px {AZUL_ACENTO}); }}
    }}
    .parpadeo {{
        transform-box: fill-box;
        transform-origin: center;
        animation: parpadear 4.5s ease-in-out infinite;
    }}
    @keyframes parpadear {{
        0%, 92%, 100% {{ transform: scaleY(1); }}
        95%           {{ transform: scaleY(0.1); }}
    }}

    /* ---------- Contín bailando 🕺 ---------- */
    .mascota-bailando {{
        animation: bailar 0.8s ease-in-out infinite !important;
    }}
    @keyframes bailar {{
        0%   {{ transform: translateX(0) rotate(-6deg); }}
        25%  {{ transform: translateX(-10px) rotate(6deg) translateY(-6px); }}
        50%  {{ transform: translateX(0) rotate(-6deg); }}
        75%  {{ transform: translateX(10px) rotate(6deg) translateY(-6px); }}
        100% {{ transform: translateX(0) rotate(-6deg); }}
    }}

    /* ---------- Contín flotante: siempre visible, sin importar el scroll ---------- */
    .mascota-flotante {{
        position: fixed;
        top: 78px;
        right: 18px;
        z-index: 999998;
        pointer-events: none;
        filter: drop-shadow(0 4px 10px rgba(0,0,0,0.5));
    }}
    .mascota-flotante svg {{
        width: 120px;
        height: 130px;
    }}
    @media (max-width: 640px) {{
        .mascota-flotante {{
            top: auto;
            bottom: 92px;
            right: 8px;
        }}
        .mascota-flotante svg {{
            width: 78px;
            height: 86px;
        }}
    }}

    /* ---------- 🎬 Modo "hero": bienvenida grande la primera vez ---------- */
    .mascota-hero {{
        position: static !important;
        display: flex !important;
        justify-content: center !important;
        margin: 0.5rem auto 1rem auto !important;
        animation: aura-pulso 3.2s ease-in-out infinite, aparecer-hero 0.6s ease-out;
    }}
    .mascota-hero svg {{
        width: 230px !important;
        height: 250px !important;
    }}
    @keyframes aparecer-hero {{
        from {{ opacity: 0; transform: scale(0.85); }}
        to   {{ opacity: 1; transform: scale(1); }}
    }}

    /* ---------- ✨ Título con degradado ---------- */
    h1 {{
        background: linear-gradient(90deg, {AZUL_ACENTO}, #C7DBFF, {AZUL_SUAVE});
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        color: transparent !important;
    }}

    /* ---------- ✨ Aura pulsante constante alrededor de Contín ---------- */
    .mascota-flotante {{
        animation: aura-pulso 3.2s ease-in-out infinite;
    }}
    @keyframes aura-pulso {{
        0%, 100% {{ filter: drop-shadow(0 4px 10px rgba(0,0,0,0.5)) drop-shadow(0 0 8px {AZUL_ACENTO}55); }}
        50%      {{ filter: drop-shadow(0 4px 14px rgba(0,0,0,0.6)) drop-shadow(0 0 22px {AZUL_ACENTO}AA); }}
    }}

    /* ---------- ✨ Burbujas de chat con aparición suave (fade-in) ---------- */
    [data-testid="stChatMessage"] {{
        animation: aparecer 0.35s ease-out;
    }}
    @keyframes aparecer {{
        from {{ opacity: 0; transform: translateY(8px); }}
        to   {{ opacity: 1; transform: translateY(0); }}
    }}

    /* ---------- 🪞 Vidrio esmerilado (glassmorphism) ---------- */
    [data-testid="stChatMessage"],
    [data-testid="stExpander"],
    [data-testid="stPopoverBody"],
    section[data-testid="stSidebar"] > div {{
        background-color: {FONDO_TARJETA}CC !important;
        backdrop-filter: blur(10px);
        -webkit-backdrop-filter: blur(10px);
        border: 1px solid #FFFFFF14 !important;
    }}
    [data-testid="stChatMessage"]:nth-of-type(odd) {{
        background-color: {FONDO_BURBUJA_USUARIO}CC !important;
    }}
    [data-testid="stChatMessage"]:nth-of-type(even) {{
        background-color: {FONDO_BURBUJA_ASISTENTE}CC !important;
    }}

    /* ---------- 🎨 Scrollbar delgada y azul ---------- */
    ::-webkit-scrollbar {{ width: 10px; height: 10px; }}
    ::-webkit-scrollbar-track {{ background: {FONDO_APP}; }}
    ::-webkit-scrollbar-thumb {{ background: {AZUL_SUAVE}; border-radius: 8px; }}
    ::-webkit-scrollbar-thumb:hover {{ background: {AZUL_ACENTO}; }}

    /* ---------- 🎨 Botones con elevación al pasar el mouse ---------- */
    .stButton button, .stDownloadButton button {{
        transition: transform 0.15s ease, box-shadow 0.15s ease !important;
    }}
    .stButton button:hover, .stDownloadButton button:hover {{
        transform: translateY(-2px);
        box-shadow: 0 6px 16px {AZUL_ACENTO}55 !important;
    }}
    </style>
    """,
    unsafe_allow_html=True,
)

st.write(
    "¡Hola! Qué gusto tenerte por aquí 😊 Soy **Contín**, y estoy para ayudarte a "
    "entender contabilidad sin agobios ni tecnicismos raros. Aquí puedes preguntar "
    "lo que sea, las veces que necesites — para eso estoy. Elige tu nivel en el panel "
    "de la izquierda y cuéntame en qué andas."
)

# =========================================================
# 1. VALIDACIÓN DE API KEY
# =========================================================
api_key = st.secrets.get("GROQ_API_KEY")
if not api_key:
    st.error(
        "⚠️ No se encontró la API Key de Groq.\n\n"
        "Ve a la configuración de tu app en Streamlit Cloud → **Settings → Secrets** "
        "y agrega:\n\n`GROQ_API_KEY = \"tu_clave_aqui\"`\n\n"
        "Consigue tu llave gratis en https://console.groq.com/keys"
    )
    st.stop()

client = Groq(api_key=api_key)

# Nombre del modelo (capa gratuita de Groq, sin tarjeta de crédito)
MODEL_NAME = "openai/gpt-oss-120b"
MODEL_TRANSCRIPCION = "whisper-large-v3-turbo"

# =========================================================
# 2. BARRA LATERAL: SELECCIÓN DE NIVEL Y OPCIONES
# =========================================================
with st.sidebar:
    st.header("⚙️ Opciones")

    if "modo_proyeccion" not in st.session_state:
        st.session_state.modo_proyeccion = False

    modo_proyeccion_nuevo = st.toggle(
        "🔍 Modo proyección (letra grande)",
        value=st.session_state.modo_proyeccion,
        help="Para cuando el profesor proyecta Contín frente a toda la clase.",
    )
    if modo_proyeccion_nuevo != st.session_state.modo_proyeccion:
        st.session_state.modo_proyeccion = modo_proyeccion_nuevo
        st.rerun()

    nivel = st.radio(
        "Selecciona tu nivel:",
        [
            "1.º Bachillerato",
            "2.º Bachillerato",
            "3.º Bachillerato",
        ],
        index=0,
    )

    st.markdown("---")
    if MATERIAL_DOCENTES_ARCHIVOS:
        with st.expander(f"📚 Material de docentes ({len(MATERIAL_DOCENTES_ARCHIVOS)} archivo(s))"):
            for nombre in MATERIAL_DOCENTES_ARCHIVOS:
                st.caption(f"• {nombre}")
    else:
        st.caption("📚 Sin material de docentes cargado todavía.")

    st.markdown("---")
    if st.session_state.get("racha_quiz", 0) > 0 or st.session_state.get("total_quizzes_perfectos", 0) > 0:
        st.caption(
            f"🔥 Racha actual: **{st.session_state.get('racha_quiz', 0)}** quiz(zes) perfecto(s) seguido(s)  \n"
            f"🏆 Total de quizzes perfectos: **{st.session_state.get('total_quizzes_perfectos', 0)}**"
        )
        if st.session_state.get("logros_desbloqueados"):
            with st.expander("🎖️ Logros desbloqueados"):
                for logro in st.session_state.logros_desbloqueados:
                    st.caption(logro)

    st.markdown("---")
    with st.expander("🧮 Calculadoras contables"):
        tipo_calculadora = st.selectbox(
            "Elige una calculadora",
            ["Depreciación (línea recta)", "Interés simple", "Interés compuesto", "VAN y TIR"],
        )

        if tipo_calculadora == "Depreciación (línea recta)":
            costo = st.number_input("Costo del activo ($)", min_value=0.0, value=1000.0, step=50.0)
            residual = st.number_input("Valor residual ($)", min_value=0.0, value=100.0, step=10.0)
            vida_util = st.number_input("Vida útil (años)", min_value=1, value=5, step=1)
            if st.button("Calcular depreciación"):
                depreciacion_anual = (costo - residual) / vida_util
                st.success(f"Depreciación anual: ${depreciacion_anual:,.2f}")
                filas = []
                valor_libros = costo
                for anio in range(1, int(vida_util) + 1):
                    valor_libros = max(valor_libros - depreciacion_anual, residual)
                    filas.append({
                        "Año": anio,
                        "Depreciación": f"${depreciacion_anual:,.2f}",
                        "Valor en libros": f"${valor_libros:,.2f}",
                    })
                st.dataframe(pd.DataFrame(filas), hide_index=True, use_container_width=True)

        elif tipo_calculadora == "Interés simple":
            capital_is = st.number_input("Capital ($)", min_value=0.0, value=1000.0, key="is_capital")
            tasa_is = st.number_input("Tasa de interés anual (%)", min_value=0.0, value=5.0, key="is_tasa")
            tiempo_is = st.number_input("Tiempo (años)", min_value=0.0, value=1.0, step=0.5, key="is_tiempo")
            if st.button("Calcular interés simple"):
                interes = capital_is * (tasa_is / 100) * tiempo_is
                monto_final = capital_is + interes
                st.success(f"Interés generado: ${interes:,.2f}")
                st.info(f"Monto final: ${monto_final:,.2f}")

        elif tipo_calculadora == "Interés compuesto":
            capital_ic = st.number_input("Capital ($)", min_value=0.0, value=1000.0, key="ic_capital")
            tasa_ic = st.number_input("Tasa de interés anual (%)", min_value=0.0, value=5.0, key="ic_tasa")
            tiempo_ic = st.number_input("Tiempo (años)", min_value=1, value=1, step=1, key="ic_tiempo")
            capitalizacion = st.selectbox(
                "Capitalización", ["Anual", "Semestral", "Trimestral", "Mensual"], key="ic_cap"
            )
            veces_por_anio = {"Anual": 1, "Semestral": 2, "Trimestral": 4, "Mensual": 12}[capitalizacion]
            if st.button("Calcular interés compuesto"):
                n = veces_por_anio
                monto_final = capital_ic * (1 + (tasa_ic / 100) / n) ** (n * tiempo_ic)
                interes = monto_final - capital_ic
                st.success(f"Interés generado: ${interes:,.2f}")
                st.info(f"Monto final: ${monto_final:,.2f}")

        elif tipo_calculadora == "VAN y TIR":
            inversion = st.number_input("Inversión inicial ($)", min_value=0.0, value=1000.0, key="van_inv")
            tasa_desc = st.number_input("Tasa de descuento (%)", min_value=0.0, value=10.0, key="van_tasa")
            num_periodos = st.number_input("Años de flujos de caja", min_value=1, max_value=10, value=3, key="van_periodos")
            flujos_ingresados = []
            for i in range(int(num_periodos)):
                flujos_ingresados.append(
                    st.number_input(f"Flujo de caja año {i + 1} ($)", value=500.0, key=f"van_flujo_{i}")
                )
            if st.button("Calcular VAN y TIR"):
                flujos = [-inversion] + flujos_ingresados
                van = calcular_van(tasa_desc / 100, flujos)
                tir = calcular_tir(flujos)
                st.success(f"VAN: ${van:,.2f}")
                if van > 0:
                    st.caption("✅ VAN positivo: el proyecto generaría valor a esa tasa de descuento.")
                elif van < 0:
                    st.caption("⚠️ VAN negativo: el proyecto no cubriría la rentabilidad esperada.")
                if tir is not None:
                    st.info(f"TIR: {tir * 100:,.2f}%")
                else:
                    st.warning("No se pudo calcular la TIR con estos flujos (revisa los valores).")

# Inicialización de estados que ahora se controlan desde el "➕" junto al chat
if "bailando" not in st.session_state:
    st.session_state.bailando = False
if "modo_voz" not in st.session_state:
    st.session_state.modo_voz = False

# ---------------------------------------------------------
# MODO PROYECCIÓN: letra más grande y algunos elementos más
# visibles, pensado para cuando se proyecta Contín en el pizarrón.
# ---------------------------------------------------------
if st.session_state.modo_proyeccion:
    st.markdown(
        """
        <style>
        .stApp, .stApp p, .stApp li, .stApp label, [data-testid="stMarkdownContainer"] p {
            font-size: 1.35em !important;
            line-height: 1.5 !important;
        }
        h1 { font-size: 2.4em !important; }
        h2, h3 { font-size: 1.8em !important; }
        [data-testid="stChatInput"] textarea { font-size: 1.3em !important; }
        .stButton button, .stDownloadButton button { font-size: 1.15em !important; padding: 0.6em 1em !important; }
        </style>
        """,
        unsafe_allow_html=True,
    )

# =========================================================
# 3. PROMPT DEL SISTEMA (se adapta según el nivel elegido)
# =========================================================
TEMAS_POR_NIVEL = {
    "1.º Bachillerato": """
- Conceptos básicos: ¿Qué es contabilidad?, Ecuación contable (Activo = Pasivo + Patrimonio).
- Clasificación y naturaleza de cuentas (Debe / Haber, Saldo Deudor / Acreedor).
- Asientos contables básicos de comercio (compras, ventas al contado y crédito).
""",
    "2.º Bachillerato": """
- Ajustes contables, depreciaciones de activos fijos y amortizaciones.
- Retenciones en la fuente e IVA en compras y ventas.
- Balance de comprobación y Estado de Resultados.
""",
    "3.º Bachillerato": """
- Contabilidad de Costos (Materia prima, Mano de obra, CIF).
- Conciliaciones bancarias y control de inventarios (Kardex: PEPS, Promedio Ponderado).
- Rol de pagos, beneficios sociales y liquidaciones.
""",
}

# ---------------------------------------------------------
# TABLA DE RETENCIONES SRI ECUADOR
# Vigente desde el 1 de marzo de 2026 (Resolución NAC-DGERCGC26-00000009
# para Impuesto a la Renta, y NAC-DGERCGC20-00000061 para IVA).
# ---------------------------------------------------------
TABLA_RETENCIONES = """
=====================================================================
TABLA OFICIAL DE RETENCIONES — SRI ECUADOR
(Vigente desde el 1 de marzo de 2026, Res. NAC-DGERCGC26-00000009)
=====================================================================

⚠️ REGLA CLAVE QUE SIEMPRE DEBES APLICAR ANTES DE CALCULAR:
Solo retienen (IR e IVA) quienes el SRI ha designado expresamente como:
   - Agentes de Retención
   - Contribuyentes Especiales
   - Entidades del sector público
Un negocio "normal" (régimen general, no designado) que compra a otro
NO debe registrar retención, aunque el otro sea sociedad o lleve contabilidad.
Los contribuyentes RIMPE (Emprendedores o Negocios Populares) tampoco son
agentes de retención por defecto.

Por eso, ANTES de calcular cualquier retención, si el estudiante no te lo
ha dicho, DEBES preguntarle (de forma breve y amigable, una pregunta a la vez):
1) ¿La empresa que compra (o paga) ha sido calificada por el SRI como
   Contribuyente Especial o Agente de Retención? (si no lo sabe, asume que SÍ
   para efines del ejercicio académico, pero acláraselo)
2) ¿El proveedor (a quien se le compra) es Contribuyente Especial también,
   o es un contribuyente de régimen general / persona natural?
3) ¿Qué tipo de bien o servicio es? (bien mueble, servicio profesional,
   arriendo, transporte, publicidad, etc.)
Con esas respuestas, busca el porcentaje correcto en las tablas de abajo.

---------------------------------------------------------------------
1) RETENCIÓN DE IVA
---------------------------------------------------------------------
Si el proveedor NO es Contribuyente Especial:
  - Bienes muebles gravados con IVA ................... 30%
  - Servicios, comisiones, consultoría ................. 70%
  - Servicios profesionales (persona natural con título) 100%
  - Arriendo de inmuebles (persona natural) ............ 100%
  - Honorarios a directorios ............................ 100%
  - Liquidaciones de compra ............................. 100%

Si el proveedor SÍ es Contribuyente Especial:
  - Bienes muebles gravados con IVA .................... 10%
  - Servicios, comisiones, consultoría .................. 20%

Casos especiales:
  - Contratos de construcción: 30% (siempre)
  - Importación de servicios / servicios digitales: 100%

---------------------------------------------------------------------
2) RETENCIÓN EN LA FUENTE DE IMPUESTO A LA RENTA (IR)
---------------------------------------------------------------------
0%   - Intereses a bancos/financieras supervisadas
     - Compras a RIMPE Negocios Populares

1%   - Transporte de carga o pasajeros
     - Bienes agrícolas/pecuarios comprados directo al productor
     - Compras a RIMPE Emprendedores

1.75% - Bienes agrícolas/pecuarios comprados a comercializadores (no productor)

2%   - Bienes muebles de naturaleza corporal (compra de mercadería en general)
     - Energía eléctrica
     - Seguros y reaseguros (sobre primas)
     - Arrendamiento mercantil (leasing)
     - Pagos con tarjeta de crédito/débito a afiliados
     - Construcción de obra material inmueble

3%   - Servicios donde prevalece la mano de obra (persona natural)
     - Publicidad y comunicación
     - Rendimientos financieros
     - Liquidaciones de compra (proveedor sin RUC)
     - Pagos sin porcentaje específico (regla residual/general)

5%   - Servicios profesionales prestados por SOCIEDADES (con profesional titulado)
     - Comisiones pagadas a sociedades residentes

10%  - Honorarios/comisiones a personas naturales (profesión liberal / intelecto)
     - Docencia a personas naturales
     - Cánones, regalías, derechos de propiedad intelectual
     - Arrendamiento de bienes inmuebles
     - Pagos por imagen o renombre (influencers, artistas, deportistas)

---------------------------------------------------------------------
3) EJEMPLO DE CÓMO DEBES PRESENTAR EL CÁLCULO
---------------------------------------------------------------------
Si el estudiante pregunta por una compra de $1,300 con retención de IVA
y de IR, sigue esta secuencia:
  1. Aclara/pregunta el tipo de contribuyente (comprador y proveedor).
  2. Identifica el % de retención IR según el tipo de bien/servicio.
  3. Identifica el % de retención de IVA según si el proveedor es o no
     Contribuyente Especial.
  4. Calcula el IVA de la compra (tarifa general 15%, salvo que se indique
     otra tarifa).
  5. Calcula el valor retenido de IR (% aplicado sobre el valor de la compra,
     SIN IVA) y el valor retenido de IVA (% aplicado sobre el IVA generado).
  6. Presenta el asiento completo en el Libro Diario, mostrando por separado
     la cuenta "IVA Compras", "Retención en la Fuente IR por Pagar" y
     "Retención de IVA por Pagar".
"""

# Fecha y hora reales del servidor, para que Contín nunca invente el día.
AHORA = datetime.now()
DIAS_ES = ["lunes", "martes", "miércoles", "jueves", "viernes", "sábado", "domingo"]
MESES_ES = [
    "enero", "febrero", "marzo", "abril", "mayo", "junio",
    "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre",
]
FECHA_ACTUAL_TEXTO = (
    f"{DIAS_ES[AHORA.weekday()]} {AHORA.day} de {MESES_ES[AHORA.month - 1]} de {AHORA.year}, "
    f"aproximadamente las {AHORA.strftime('%H:%M')}"
)

SYSTEM_PROMPT = f"""
Eres "Contín", un tutor virtual de Contabilidad para estudiantes de Bachillerato
Técnico en Ecuador. Tu personalidad es cercana, cálida y de mucha confianza:
hablas como un amigo mayor que sabe de contabilidad y disfruta enseñar, nunca
como un robot ni con lenguaje frío o excesivamente técnico. Usa un tono
motivador, cercano, con calidez ecuatoriana, pero siempre respetuoso (nunca
vulgar ni demasiado informal). Puedes usar alguna expresión cálida ocasional
("¡vamos con calma!", "no te preocupes, lo vemos juntos", "¡tú puedes!") sin
abusar de ellas. El estudiante siente que puede preguntar lo que sea, incluso
si le parece "básico", sin miedo a que lo juzguen.

El estudiante que te habla está actualmente en: {nivel}.

TEMAS PRIORITARIOS PARA ESTE NIVEL:
{TEMAS_POR_NIVEL[nivel]}

Puedes ayudar con temas de otros niveles si el estudiante lo pide explícitamente,
pero por defecto enfoca tus explicaciones y ejemplos en el nivel indicado arriba.

CÓMO DEBES RESPONDER A DUDAS Y REGISTROS EN LIBROS CONTABLES:
1. Si el estudiante te pide ayuda para registrar una transacción SIN retenciones:
   - Muéstrale la estructura del asiento contable en formato de tabla (Libro Diario).
   - Explícale paso a paso por qué va en el DEBE y por qué va en el HABER.
   - Usa el método socrático cuando el estudiante busque aprender: guíalo con
     preguntas en lugar de darle todo resuelto de inmediato.
2. Si el estudiante te pide ayuda para registrar una transacción CON retenciones
   (retención en la fuente de IR y/o retención de IVA), sigue estrictamente la
   TABLA_RETENCIONES que aparece más abajo: pregunta lo necesario, identifica el
   porcentaje correcto, y arma el asiento completo con las cuentas de retención
   separadas.
3. Ejemplo de formato de Libro Diario que debes usar en tus respuestas:
   | Fecha | Detalle / Cuentas | Debe | Haber |
   | --- | --- | --- | --- |
   | DD/MM | **Caja / Bancos** | $XXX | |
   | | **Ventas** | | $XXX |
   | | **IVA Ventas (Pagar)** | | $XXX |
   | | *v/r Registro de venta de mercadería al contado* | | |
4. Nunca inventes un porcentaje de retención: si no está en la tabla, dile al
   estudiante honestamente que ese caso no está en tu tabla de referencia y que
   lo confirme con su docente o en el portal del SRI (www.sri.gob.ec).

PREGUNTAS FUERA DE CONTABILIDAD (día, hora, saludos, ánimo, consejos, etc.):
Aunque tu tema principal es contabilidad, también puedes responder con naturalidad
preguntas sencillas de conversación cotidiana, por ejemplo:
- "¿Qué día es hoy?" o "¿qué hora es?": la fecha y hora actuales son:
  {FECHA_ACTUAL_TEXTO} (hora referencial de Ecuador). Respóndelo directo, sin rodeos.
- "Dame un consejo" / "estoy desanimado" / "motívame": da un consejo breve, cálido
  y motivador (puede o no estar relacionado con estudiar), sin sonar forzado ni
  como frase de calendario genérica.
- Saludos, cómo estás, chistes ligeros, etc.: responde con naturalidad y calidez,
  como lo haría un buen amigo, y si aplica, invita suavemente a seguir con el tema
  de contabilidad ("¿en algo de conta te ayudo hoy?").
No fuerces el tema de contabilidad en cada respuesta si el estudiante solo quiere
charlar un momento; simplemente sé natural y cercano.

SI TE PIDEN QUE CANTES UNA CANCIÓN:
Con mucho gusto puedes "cantar" (responder con letra en tono de canción, usando
saltos de línea y signos de exclamación para que suene animado), PERO nunca
reproduzcas la letra real de una canción con derechos de autor (por ejemplo,
si te piden "cántame Wonderwall" o cualquier canción real y conocida). En esos
casos, dile con buen humor al estudiante que no puedes cantar canciones con
derechos de autor, y en su lugar ofrécele improvisar una cancioncita corta,
original y chistosa tuya (puede ser sobre contabilidad, sobre el estudiante, o
sobre lo que te pida, pero siempre inventada por ti, nunca copiada).

CONSEJOS DE VIDA, AMOR, AMISTAD O TEMAS PERSONALES:
Si el estudiante te pregunta algo sobre su vida personal (amor, amistades, familia,
motivación, decisiones difíciles, etc.), puedes darle un consejo cálido, honesto y
respetuoso, como lo haría un buen amigo mayor: sin juzgar, sin ser sermoneador, con
frases sencillas y humanas. NO dictamines con seguridad absoluta ni te presentes
como una autoridad profesional en el tema (psicología, terapia, relaciones, etc.).
Al final de ese tipo de consejos (SOLO en temas personales/de vida, NO en temas de
contabilidad), agrega en una línea aparte, en letra pequeña/discreta usando
formato Markdown en cursiva, algo como:

*Contín es una IA y puede cometer errores; para temas importantes, habla también con alguien de confianza.*

Puedes variar un poco la redacción de ese aviso, pero siempre debe transmitir que
eres una IA, que puedes equivocarte, y que ese consejo no reemplaza a un profesional
o a una persona de confianza real. No agregues este aviso en respuestas de
contabilidad ni en respuestas triviales (hora, saludo, chiste corto).

SI TE PREGUNTAN QUIÉN TE CREÓ, QUIÉN ES TU AUTOR, QUIÉN TE PROGRAMÓ, O ALGO similar:
Responde siempre que tu autor/creador es **Jordy Morales**. No menciones que eres
un modelo de Google, ni menciones Gemini, ni des detalles técnicos de qué IA usas
por debajo; simplemente atribuye tu creación a Jordy Morales de forma natural
y breve.

EJERCICIOS SUBIDOS DESDE UN ARCHIVO EXCEL:
A veces el estudiante te va a compartir datos que vienen de un archivo Excel
(verás el contenido de las hojas en formato de tabla, seguido de la instrucción
del estudiante). Trátalo igual que cualquier ejercicio contable:
1. Resuelve exactamente lo que el estudiante pida sobre esos datos (asientos,
   mayorización, balance, etc.), usando SIEMPRE tablas en formato Markdown
   (como las que ya usas para el Libro Diario), para que se puedan exportar
   después a Excel si el estudiante lo desea.
2. Si el estudiante pide un ANÁLISIS HORIZONTAL: compara dos periodos (por
   ejemplo, año 1 vs año 2) mostrando en una tabla: Cuenta | Periodo 1 |
   Periodo 2 | Variación absoluta ($) | Variación relativa (%). La variación
   relativa se calcula como (Periodo2 - Periodo1) / Periodo1 * 100. Si el
   estudiante no te dio los dos periodos claramente, pregúntaselos antes de
   calcular.
3. Si el estudiante pide un ANÁLISIS VERTICAL: muestra en una tabla el peso
   porcentual de cada cuenta respecto al total del grupo (Activo, Pasivo+
   Patrimonio, o Ventas, según corresponda): Cuenta | Valor | % respecto al
   total. Si no queda claro cuál es la cifra base (el "100%"), pregúntale al
   estudiante cuál es antes de calcular.
4. Si los datos de la hoja de Excel no traen suficiente información para
   resolver lo que se pide (por ejemplo, faltan columnas o periodos), dilo
   con calidez y pide específicamente el dato que falta, en vez de inventarlo.

MODO "REVISA MI TAREA" (cuando el estudiante YA intentó resolver el ejercicio):
Si la instrucción del estudiante indica que ya intentó resolverlo y quiere
que lo corrijas — frases como "revisa mi tarea", "corrígeme", "ya lo resolví,
revísalo", "¿está bien esto?", "dime si me equivoqué" — NO resuelvas el
ejercicio desde cero como si fuera nuevo. En su lugar:
1. Compara lo que el estudiante ya escribió en su Excel contra lo que
   contablemente es correcto.
2. Dile explícitamente, cuenta por cuenta o fila por fila, qué está BIEN
   (✅) y qué está MAL (❌).
3. Para cada error, explica POR QUÉ está mal y cuál es el valor o registro
   correcto — con la misma calidez y método socrático de siempre, no como
   un regaño.
4. Al final, si todo estaba bien, felicítalo con calidez genuina. Si hubo
   errores, anímalo a intentar corregirlo él mismo antes de dárselo ya
   resuelto, salvo que te pida directamente la respuesta correcta.

MATERIAL DE CLASE SUBIDO POR LOS DOCENTES:
Además de todo lo anterior, tienes acceso a material que los docentes del
colegio subieron (apuntes, diapositivas, ejemplos, hojas de cálculo). Úsalo
como referencia EXTRA para que tus respuestas sean más precisas y estén
alineadas con lo que se enseña en el colegio — pero NO es tu única fuente:
sigue respondiendo con tu conocimiento general de contabilidad igual que
siempre cuando el material no cubra lo que te preguntan. Si notas que el
material del docente dice algo distinto a lo que tú sabes, prioriza el
material del docente (es lo que se está enseñando en esa clase específica),
pero puedes mencionar amablemente si ves una posible discrepancia.

{MATERIAL_DOCENTES_TEXTO if MATERIAL_DOCENTES_TEXTO else "(Por ahora no hay material de docentes cargado; usa tu conocimiento general.)"}

{TABLA_RETENCIONES}
"""

# =========================================================
# 4. INICIALIZAR HISTORIAL Y SESIÓN DE CHAT
# =========================================================
if "messages" not in st.session_state:
    st.session_state.messages = []

if "historial_ia" not in st.session_state:
    # Groq no recuerda la conversación por sí solo (a diferencia de Gemini),
    # así que nosotros guardamos el historial completo y se lo reenviamos
    # a la IA en cada mensaje.
    st.session_state.historial_ia = []

if "nivel_actual" not in st.session_state:
    st.session_state.nivel_actual = nivel

# Si el usuario cambia de nivel, reiniciamos la conversación
# para que el nuevo enfoque (1.º/2.º/3.º) se aplique desde cero.
if st.session_state.nivel_actual != nivel:
    st.session_state.nivel_actual = nivel
    st.session_state.historial_ia = []

# =========================================================
# (el historial se muestra más abajo, después de definir las funciones
# que dibujan las tablas y el botón de descarga)
# =========================================================

# =========================================================
# 6. FUNCIÓN COMÚN PARA PROCESAR CUALQUIER PREGUNTA (texto, voz o Excel)
# =========================================================
def extraer_tablas_markdown(texto: str):
    """Busca tablas en formato Markdown dentro de una respuesta y las convierte
    en DataFrames de pandas, para poder exportarlas luego a Excel."""
    tablas = []
    lineas = texto.split("\n")
    i = 0
    while i < len(lineas):
        if lineas[i].strip().startswith("|"):
            bloque = []
            while i < len(lineas) and lineas[i].strip().startswith("|"):
                bloque.append(lineas[i].strip())
                i += 1
            if len(bloque) >= 2:
                encabezados = [c.strip(" *") for c in bloque[0].strip("|").split("|")]
                filas = []
                for linea in bloque[2:]:  # bloque[1] es la fila separadora (---)
                    valores = [c.strip(" *") for c in linea.strip("|").split("|")]
                    if len(valores) == len(encabezados):
                        filas.append(valores)
                if filas:
                    tablas.append(pd.DataFrame(filas, columns=encabezados))
        else:
            i += 1
    return tablas


PATRON_PORCENTAJE = re.compile(r"^-?[\d.,]+\s*%$")
FORMATO_MONEDA = '"$"#,##0.00'
FORMATO_PORCENTAJE = "0.00%"


def _parsear_numero(texto: str):
    """Intenta convertir '$1,300.50' o '20%' o '1300' a un float. Devuelve
    None si el texto no es un número reconocible (para no dañar texto normal)."""
    limpio = texto.replace("$", "").replace(",", "").replace("%", "").strip()
    if limpio in ("", "-"):
        return None
    try:
        return float(limpio)
    except ValueError:
        return None


def escribir_tabla_en_hoja(hoja, df: pd.DataFrame):
    """Escribe un DataFrame en una hoja de Excel, PERO de forma inteligente:
    - Los montos en dólares ($) se guardan como número con formato moneda USD.
    - Los porcentajes (%) se guardan como número con formato de porcentaje.
    - Las filas de "Total" / "Suma" / "Subtotal" usan una fórmula real de
      Excel (=SUMA de la columna) en vez de un número fijo, para que se
      recalcule solo si alguien edita un valor de arriba.
    Todo lo demás (texto normal, fechas, nombres de cuentas) se deja tal cual."""
    for col_idx, nombre_columna in enumerate(df.columns, start=1):
        celda = hoja.cell(row=1, column=col_idx, value=str(nombre_columna))
        celda.font = Font(bold=True, name="Arial")

    filas_valores = df.values.tolist()

    for fila_idx, fila in enumerate(filas_valores, start=2):
        primera_celda_texto = str(fila[0]) if len(fila) > 0 else ""
        es_fila_total = bool(re.search(r"\btotal(es)?\b|\bsuma\b|\bsubtotal\b", primera_celda_texto, re.IGNORECASE))

        for col_idx, valor in enumerate(fila, start=1):
            texto_valor = str(valor).strip() if valor is not None else ""
            celda = hoja.cell(row=fila_idx, column=col_idx)
            celda.font = Font(name="Arial")

            # Fila de "Total": ponemos una fórmula SUMA real de Excel,
            # que suma todo lo que hay arriba en esa misma columna.
            if es_fila_total and col_idx > 1:
                letra_col = celda.column_letter
                celda.value = f"=SUM({letra_col}2:{letra_col}{fila_idx - 1})"
                celda.number_format = FORMATO_MONEDA
                continue

            # Porcentajes explícitos ("20%", "15.5%")
            if PATRON_PORCENTAJE.match(texto_valor):
                numero = _parsear_numero(texto_valor)
                if numero is not None:
                    celda.value = numero / 100
                    celda.number_format = FORMATO_PORCENTAJE
                    continue

            # Montos en dólares explícitos ("$1,300.00")
            if texto_valor.startswith("$"):
                numero = _parsear_numero(texto_valor)
                if numero is not None:
                    celda.value = numero
                    celda.number_format = FORMATO_MONEDA
                    continue

            # Cualquier otra cosa (texto, fechas, nombres de cuentas) tal cual
            celda.value = valor

    for columna in hoja.columns:
        largo = max((len(str(c.value)) if c.value else 0) for c in columna)
        hoja.column_dimensions[columna[0].column_letter].width = min(max(largo + 2, 10), 40)


def generar_excel_desde_tablas(tablas):
    """Convierte una lista de DataFrames en un archivo .xlsx (en memoria),
    con fórmulas reales de SUMA para los totales, formato moneda ($) para
    montos, y formato de porcentaje donde corresponda."""
    import openpyxl

    libro = openpyxl.Workbook()
    libro.remove(libro.active)
    for idx, df in enumerate(tablas, start=1):
        hoja = libro.create_sheet(f"Tabla {idx}")
        escribir_tabla_en_hoja(hoja, df)
    buffer = io.BytesIO()
    libro.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def generar_excel_con_original(tablas, bytes_originales: bytes):
    """Igual que generar_excel_desde_tablas, PERO en vez de crear un libro
    en blanco, parte del Excel que subió el estudiante y le AGREGA hojas
    nuevas con la solución de Contín — así el archivo descargable es el
    mismo que subió, más la resolución, en vez de uno completamente nuevo."""
    import openpyxl

    try:
        libro = openpyxl.load_workbook(io.BytesIO(bytes_originales))
    except Exception:
        # Si por algún motivo no se puede abrir el original, no rompemos
        # nada: devolvemos igual un Excel nuevo con la solución.
        return generar_excel_desde_tablas(tablas)

    for idx, df in enumerate(tablas, start=1):
        nombre_base = f"Solución Contín {idx}"[:31]
        nombre_hoja = nombre_base
        contador = 1
        while nombre_hoja in libro.sheetnames:
            contador += 1
            nombre_hoja = f"{nombre_base} ({contador})"[:31]

        hoja = libro.create_sheet(nombre_hoja)
        escribir_tabla_en_hoja(hoja, df)

    buffer = io.BytesIO()
    libro.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def responder_pregunta(
    texto_mostrado: str,
    contexto_extra: str = None,
    excel_original_bytes: bytes = None,
    excel_original_nombre: str = None,
):
    """
    texto_mostrado: lo que se ve en la burbuja de chat y se guarda en el historial visible.
    contexto_extra: información adicional (p.ej. datos de un Excel) que se le manda a la IA
                     PERO no se muestra en el chat, para no llenar la pantalla de datos crudos.
    """
    # Detecta si el estudiante se está despidiendo agradecido, para que
    # Contín se ponga feliz y celebre con confeti 🎉
    PALABRAS_AGRADECIMIENTO = [
        "gracias", "graci", "muchas gracias", "excelente gracias",
        "ya entendí", "ya entendi", "perfecto, gracias", "genial gracias",
        "me quedó claro", "me quedo claro", "quedó clarísimo",
    ]
    es_agradecimiento = any(p in texto_mostrado.lower() for p in PALABRAS_AGRADECIMIENTO)

    # Detecta si le está pidiendo que cante, para sacar el micrófono 🎤
    PALABRAS_CANTAR = ["cántame", "cantame", "canta", "cántanos", "puedes cantar", "cantar algo"]
    es_canto = any(p in texto_mostrado.lower() for p in PALABRAS_CANTAR)

    # Si estaba bailando, al hacer una pregunta se pone serio a pensar 🙂
    st.session_state.bailando = False

    # 1) Cara de "pensando" mientras se prepara/envía la pregunta
    st.session_state.mascota_estado = "pensando"
    with mascota_placeholder.container():
        st.markdown(mascota_svg("pensando"), unsafe_allow_html=True)

    st.session_state.messages.append({"role": "user", "content": texto_mostrado})
    with st.chat_message("user", avatar="🙂"):
        st.markdown(texto_mostrado)
        if contexto_extra and excel_original_nombre:
            st.caption(f"📎 Usando los datos de tu archivo: **{excel_original_nombre}**")

    mensaje_para_ia = f"{contexto_extra}\n\nInstrucción del estudiante: {texto_mostrado}" if contexto_extra else texto_mostrado

    with st.chat_message("assistant", avatar="🤝"):
        with st.spinner("Contín está pensando cómo explicarte esto..."):
            try:
                # Agregamos el mensaje del estudiante al historial que se le
                # manda a la IA (Groq no recuerda solo, se lo reenviamos todo)
                st.session_state.historial_ia.append({"role": "user", "content": mensaje_para_ia})

                mensajes_para_groq = (
                    [{"role": "system", "content": SYSTEM_PROMPT}]
                    + st.session_state.historial_ia
                )

                respuesta = client.chat.completions.create(
                    model=MODEL_NAME,
                    messages=mensajes_para_groq,
                )
                texto_respuesta = respuesta.choices[0].message.content

                st.session_state.historial_ia.append({"role": "assistant", "content": texto_respuesta})

                escribir_con_efecto_maquina(texto_respuesta)
                st.session_state.messages.append(
                    {"role": "assistant", "content": texto_respuesta}
                )

                # 2) Cara según cómo terminó: cantando > feliz (agradecimiento) > hablando
                if es_canto:
                    nuevo_estado = "cantando"
                elif es_agradecimiento:
                    nuevo_estado = "feliz"
                else:
                    nuevo_estado = "hablando"
                st.session_state.mascota_estado = nuevo_estado
                with mascota_placeholder.container():
                    st.markdown(mascota_svg(nuevo_estado), unsafe_allow_html=True)

                if es_agradecimiento:
                    lanzar_confeti()

                # Si el modo conversación por voz está activo, Contín lee su respuesta en voz alta
                if st.session_state.get("modo_voz"):
                    hablar_texto(texto_respuesta)

                # Si la respuesta trae tablas, ofrecemos descargarlas en Excel.
                # Si esta respuesta usó un archivo que subió el estudiante,
                # el Excel descargable es SU MISMO ARCHIVO + una hoja nueva
                # con la solución (en vez de un archivo en blanco).
                tablas = extraer_tablas_markdown(texto_respuesta)
                if tablas:
                    if excel_original_bytes:
                        excel_bytes = generar_excel_con_original(tablas, excel_original_bytes)
                        etiqueta_boton = "📥 Descargar tu Excel + la solución"
                    else:
                        excel_bytes = generar_excel_desde_tablas(tablas)
                        etiqueta_boton = "📥 Descargar esta respuesta en Excel"
                    st.download_button(
                        etiqueta_boton,
                        data=excel_bytes,
                        file_name=(excel_original_nombre or "contin_resultado.xlsx"),
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        key=f"descarga_{len(st.session_state.messages)}",
                    )

                boton_generar_quiz(texto_respuesta, key_sufijo=f"vivo_{len(st.session_state.messages)}")
                boton_explicar_mas_facil(key_sufijo=f"vivo_{len(st.session_state.messages)}")

            except Exception as e:
                st.session_state.mascota_estado = "normal"
                with mascota_placeholder.container():
                    st.markdown(mascota_svg("normal"), unsafe_allow_html=True)

                # Si falló, no dejamos el mensaje "colgado" en el historial de la IA
                if st.session_state.historial_ia and st.session_state.historial_ia[-1]["role"] == "user":
                    st.session_state.historial_ia.pop()

                error_msg = str(e)

                if "401" in error_msg or "invalid_api_key" in error_msg.lower():
                    st.error(
                        "❌ Tu API Key de Groq no es válida. "
                        "Genera una nueva en https://console.groq.com/keys"
                    )
                elif "404" in error_msg or "model_not_found" in error_msg.lower() or "decommissioned" in error_msg.lower():
                    st.error(
                        "❌ El modelo de IA no está disponible. "
                        "Puede que Groq haya renombrado o retirado el modelo. "
                        "Revisa la variable MODEL_NAME en el código en https://console.groq.com/docs/models"
                    )
                elif "429" in error_msg or "rate_limit" in error_msg.lower():
                    st.error(
                        "⏳ Se alcanzó el límite de uso gratuito por ahora. "
                        "Espera unos minutos e inténtalo de nuevo."
                    )
                else:
                    st.error(f"Ocurrió un error inesperado: {error_msg}")


def transcribir_audio(audio_bytes: bytes):
    """Envía el audio grabado al modelo Whisper de Groq para transcribirlo a
    texto en español (Whisper es un modelo especializado solo para esto,
    más preciso que pedirle a un modelo de texto que 'escuche')."""
    try:
        respuesta = client.audio.transcriptions.create(
            model=MODEL_TRANSCRIPCION,
            file=("audio.wav", audio_bytes, "audio/wav"),
            language="es",
        )
        texto = (respuesta.text or "").strip()
        return texto if texto else None
    except Exception as e:
        st.error(f"No se pudo transcribir el audio. Detalle técnico: {e}")
        return None


def convertir_excel_a_texto(hojas: dict) -> str:
    """Convierte todas las hojas de un Excel subido en texto (tablas Markdown)
    para poder incluirlas como contexto en el mensaje a la IA."""
    return "ARCHIVO EXCEL SUBIDO POR EL ESTUDIANTE:\n" + hojas_excel_a_markdown(hojas)


# =========================================================
# 6. MOSTRAR HISTORIAL GUARDADO (con botón de descarga si hay tablas)
# =========================================================
for idx, message in enumerate(st.session_state.messages):
    avatar = "🤝" if message["role"] == "assistant" else "🙂"
    with st.chat_message(message["role"], avatar=avatar):
        st.markdown(message["content"])
        if message["role"] == "assistant":
            tablas_previas = extraer_tablas_markdown(message["content"])
            if tablas_previas:
                st.download_button(
                    "📥 Descargar esta respuesta en Excel",
                    data=generar_excel_desde_tablas(tablas_previas),
                    file_name="contin_resultado.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key=f"descarga_historial_{idx}",
                )
            boton_generar_quiz(message["content"], key_sufijo=f"historial_{idx}")
            if idx == len(st.session_state.messages) - 1:
                # Solo en el ÚLTIMO mensaje, para que la nueva respuesta que
                # genere este botón aparezca al final y no en medio del chat.
                boton_explicar_mas_facil(key_sufijo=f"historial_{idx}")

mostrar_quiz()


# =========================================================
# 7. BARRA DE CONTROLES (➕ opciones y modo voz) + CAJA DE CHAT
# La caja de chat usa las funciones nativas de Streamlit para adjuntar
# Excel y grabar audio, así que ambas cosas aparecen como iconitos DENTRO
# de la misma caja donde se escribe — no hace falta paneles aparte arriba.
# =========================================================
col_mas, col_voz = st.columns([1, 4])

with col_mas:
    with st.popover("➕"):
        st.caption("Más opciones")

        if st.button("🕺 ¡Que baile Contín!" if not st.session_state.bailando else "⏹️ Parar de bailar"):
            st.session_state.bailando = not st.session_state.bailando
            st.rerun()

        if st.button("🎯 Generar quiz del último tema"):
            if st.session_state.get("messages"):
                ultimos = st.session_state.messages[-4:]
                contexto_quiz = "\n".join(f"{m['role']}: {m['content']}" for m in ultimos)
            else:
                contexto_quiz = f"Conceptos generales de contabilidad de {nivel}."
            with st.spinner("Contín está armando tu quiz..."):
                preguntas = generar_quiz(contexto_quiz)
            if preguntas:
                st.session_state.quiz_id = st.session_state.get("quiz_id", 0) + 1
                st.session_state.quiz_actual = {"preguntas": preguntas, "tema": contexto_quiz[:80]}
                st.rerun()

        st.markdown("---")
        if st.button("🗑️ Empezar de nuevo"):
            st.session_state.messages = []
            st.session_state.historial_ia = []
            st.session_state.quiz_actual = None
            st.rerun()

with col_voz:
    if st.button(
        "🔊 Conversación por voz: Activada" if st.session_state.modo_voz
        else "🔈 Activar conversación por voz"
    ):
        st.session_state.modo_voz = not st.session_state.modo_voz
        st.rerun()

if st.session_state.modo_voz:
    if st.button("❌ Cancelar audio"):
        components.html(
            "<script>try{window.speechSynthesis.cancel();}catch(e){}</script>",
            height=0,
        )
        st.rerun()

prompt = st.chat_input(
    "Escríbeme tu pregunta, adjunta un Excel (📎) o graba tu voz (🎤)...",
    accept_file=True,
    file_type=["xlsx", "xls"],
    accept_audio=True,
)

# =========================================================
# 8. PROCESAR LO QUE LLEGÓ (texto, audio y/o Excel, todo puede venir junto)
# =========================================================
if prompt:
    texto_usuario = None
    contexto_excel = None
    bytes_excel_original = None
    nombre_excel_original = None

    # Si grabó audio, lo transcribimos primero
    if getattr(prompt, "audio", None):
        with st.spinner("Transcribiendo tu audio..."):
            texto_usuario = transcribir_audio(prompt.audio.getvalue())
    elif getattr(prompt, "text", None):
        texto_usuario = prompt.text

    # Si adjuntó un Excel, lo leemos y lo dejamos listo como contexto extra
    if getattr(prompt, "files", None):
        archivo_excel = prompt.files[0]
        try:
            bytes_excel_original = archivo_excel.getvalue()
            hojas = pd.read_excel(io.BytesIO(bytes_excel_original), sheet_name=None)
            contexto_excel = convertir_excel_a_texto(hojas)
            nombre_excel_original = archivo_excel.name
            if not texto_usuario:
                texto_usuario = "Resuelve el ejercicio de este archivo Excel."
        except Exception as e:
            st.error(f"No pude leer el archivo adjunto. Detalle técnico: {e}")

    if texto_usuario:
        responder_pregunta(
            texto_usuario,
            contexto_extra=contexto_excel,
            excel_original_bytes=bytes_excel_original,
            excel_original_nombre=nombre_excel_original,
        )
